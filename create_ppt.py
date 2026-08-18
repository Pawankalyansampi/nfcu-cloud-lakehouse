"""NFCU Cloud lakehouse PPT — business problem, architecture, people."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0A, 0x25, 0x40)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF3, 0xEB)
SLATE = RGBColor(0x3D, 0x4A, 0x57)
LINE = RGBColor(0xD8, 0xD2, 0xC4)

W = Inches(13.333)
H = Inches(7.5)
OUT = Path(__file__).resolve().parent / "presentation" / "NFCU_Cloud_Project.pptx"


def box(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def tb(slide, l, t, w, h, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(l, t, w, h)
    s.text_frame.word_wrap = True
    try:
        s.text_frame._txBody.bodyPr.set(
            "anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor]
        )
    except Exception:
        pass
    p = s.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return s


def footer(slide, page, total=7):
    box(slide, Inches(0), Inches(7.15), W, Inches(0.35), NAVY)
    tb(
        slide,
        Inches(0.4),
        Inches(7.16),
        Inches(10),
        Inches(0.32),
        "Navy Federal Credit Union  |  Vienna, VA  |  Cloud Lakehouse  |  Apr 2025 – Present",
        12,
        False,
        GOLD,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    tb(
        slide,
        Inches(11.4),
        Inches(7.16),
        Inches(1.6),
        Inches(0.32),
        f"{page} / {total}",
        12,
        False,
        GOLD,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def header(slide, title):
    box(slide, Inches(0), Inches(0), W, Inches(1.15), NAVY)
    box(slide, Inches(0), Inches(1.15), W, Inches(0.08), GOLD)
    tb(slide, Inches(0.5), Inches(0.28), Inches(12), Inches(0.7), title, 28, True, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def card(slide, l, t, w, h, title, body):
    box(slide, l, t, w, h, WHITE, LINE)
    box(slide, l, t, Inches(0.1), h, GOLD)
    tb(slide, l + Inches(0.3), t + Inches(0.15), w - Inches(0.45), Inches(0.4), title, 18, True, NAVY)
    tb(slide, l + Inches(0.3), t + Inches(0.6), w - Inches(0.45), h - Inches(0.75), body, 15, False, SLATE)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, NAVY)
    box(s, Inches(0), Inches(0), Inches(0.2), H, GOLD)
    tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.35), "NAVY FEDERAL CREDIT UNION", 14, True, GOLD)
    tb(s, Inches(0.7), Inches(1.85), Inches(12), Inches(1.0), "Cloud Financial Data Lakehouse", 36, True, WHITE)
    tb(
        s,
        Inches(0.7),
        Inches(2.95),
        Inches(12),
        Inches(0.5),
        "Business problem, AWS architecture, and the people who delivered it",
        20,
        False,
        GOLD,
    )
    tb(
        s,
        Inches(0.7),
        Inches(4.0),
        Inches(12),
        Inches(1.3),
        "Vienna, VA  ·  Apr 2025 – Present\nSenior Data Engineer  ·  Badhari Swaroop",
        18,
        False,
        WHITE,
    )
    tb(
        s,
        Inches(0.7),
        Inches(6.3),
        Inches(12),
        Inches(0.4),
        "Terraform  ·  S3  ·  Kinesis  ·  Glue  ·  Athena  ·  Databricks  ·  Docker  ·  Streamlit",
        16,
        False,
        GOLD,
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, CREAM)
    header(s, "Business problem")
    card(
        s,
        Inches(0.45),
        Inches(1.5),
        Inches(6.15),
        Inches(5.3),
        "What Navy Federal needed",
        "Member payments and bank-account activity sat in separate systems. Fraud review, finance reporting, and policy questions each used a different extract.\n\n"
        "That caused late fraud flags, mismatched daily totals, and no single place to query governed lake data.\n\n"
        "The credit union needed one cloud lakehouse: batch + near-real-time events, quality checks, and self-serve SQL for reporting.",
    )
    card(
        s,
        Inches(6.8),
        Inches(1.5),
        Inches(6.1),
        Inches(5.3),
        "What we delivered",
        "• One S3 lake (bronze / silver / gold) for Plaid-style accounts and PaySim payments\n"
        "• Kinesis for payment events, landed to S3\n"
        "• Databricks PySpark batch + Structured Streaming\n"
        "• Glue Catalog + Athena for governed SQL\n"
        "• FastAPI + Streamlit for APIs and dashboards\n"
        "• Terraform so the environment can be created and torn down for cost control",
    )
    footer(s, 2)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, CREAM)
    header(s, "People who worked on this")
    people = [
        ("Badhari Swaroop", "Senior Data Engineer", "Owned Terraform, the Python lakehouse, Kinesis producer, APIs, Docker, and Databricks notebooks."),
        ("Data Architect", "Platform design", "Set the medallion lake (bronze / silver / gold), Glue catalog, and Athena workgroup pattern."),
        ("Fraud Operations", "Business consumer", "Uses fraud flags and gold fraud summary for same-day payment review."),
        ("Finance / Reporting", "Business consumer", "Uses Athena gold volume and Streamlit KPIs as the daily numbers."),
        ("Compliance / Privacy", "Control partner", "Keeps secrets out of git, no full account numbers on the dashboard, demo policies in RAG."),
        ("Product Owner", "Priority and acceptance", "Accepted the demo: lake files on S3, Kinesis events, Databricks gold counts, Athena SQL."),
    ]
    for i, (name, role, work) in enumerate(people):
        col, row = i % 2, i // 2
        x = Inches(0.45) + col * Inches(6.4)
        y = Inches(1.5) + row * Inches(1.75)
        box(s, x, y, Inches(6.15), Inches(1.55), WHITE, LINE)
        tb(s, x + Inches(0.25), y + Inches(0.12), Inches(5.7), Inches(0.35), name, 16, True, NAVY)
        tb(s, x + Inches(0.25), y + Inches(0.45), Inches(5.7), Inches(0.28), role, 13, True, GOLD)
        tb(s, x + Inches(0.25), y + Inches(0.75), Inches(5.7), Inches(0.65), work, 13, False, SLATE)
    footer(s, 3)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, CREAM)
    header(s, "Cloud architecture")
    steps = [
        ("1. Sources", "Plaid-style accounts\nPaySim payments"),
        ("2. Ingest", "Batch → S3 bronze\nEvents → Kinesis"),
        ("3. Process", "Databricks PySpark\nStructured Streaming"),
        ("4. Lake", "S3 + Glue Catalog\nbronze / silver / gold"),
        ("5. Serve", "Athena SQL\nFastAPI + Streamlit"),
    ]
    for i, (title, body) in enumerate(steps):
        x = Inches(0.35) + i * Inches(2.58)
        box(s, x, Inches(1.55), Inches(2.42), Inches(2.85), WHITE, LINE)
        box(s, x, Inches(1.55), Inches(2.42), Inches(0.5), NAVY)
        tb(s, x, Inches(1.55), Inches(2.42), Inches(0.5), title, 15, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tb(s, x + Inches(0.08), Inches(2.15), Inches(2.26), Inches(2.1), body, 14, False, SLATE, PP_ALIGN.CENTER)
        if i < 4:
            tb(s, x + Inches(2.28), Inches(2.7), Inches(0.35), Inches(0.4), "→", 20, True, GOLD, PP_ALIGN.CENTER)
    box(s, Inches(0.35), Inches(4.6), Inches(12.6), Inches(2.25), WHITE, LINE)
    tb(s, Inches(0.55), Inches(4.72), Inches(12.2), Inches(0.35), "Platform (what we actually ran)", 16, True, NAVY)
    tb(
        s,
        Inches(0.55),
        Inches(5.15),
        Inches(12.2),
        Inches(1.5),
        "Terraform provisions S3, Glue, Athena, Kinesis, and CloudWatch. Docker runs API, Streamlit, and Airflow. "
        "GitHub Actions validates Terraform. Databricks Free Edition runs Spark (serverless). "
        "Quality checks are Python (Great Expectations-style). Dashboard is Streamlit, not Tableau. "
        "RDS, Redshift, VPC, and KMS were removed to control student cost.",
        14,
        False,
        SLATE,
    )
    footer(s, 4)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, CREAM)
    header(s, "What we implemented vs the target diagram")
    rows = [
        ("Ran for real", "Terraform, S3 lake, Kinesis (13 events), Glue, Athena, Databricks batch (gold 5/2/3), streaming notebook, Docker, Streamlit, CloudWatch"),
        ("Student stand-in", "Airflow in Docker (not MWAA). dbt as gold SQL files. Quality in Python. Streamlit instead of Tableau. Streaming uses availableNow on serverless (not live Kinesis→Spark)"),
        ("Not in this demo", "Snowflake trial (code only). Tableau. Great Expectations product. VPC / KMS / RDS / Redshift (removed for cost)"),
    ]
    for i, (title, body) in enumerate(rows):
        card(s, Inches(0.45), Inches(1.5) + i * Inches(1.75), Inches(12.4), Inches(1.6), title, body)
    footer(s, 5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, CREAM)
    header(s, "Business outcomes")
    outcomes = [
        ("One lake for payments and accounts", "Bronze / silver / gold on S3, cataloged in Glue, queried in Athena."),
        ("Faster fraud review", "High-risk TRANSFER and CASH_OUT land in gold fraud summary (2 types in Databricks)."),
        ("Near-real-time path", "Kinesis → S3 landing. Databricks Structured Streaming for event processing."),
        ("Cost control", "No Redshift. Tear down with terraform destroy. Typical demo day well under $1."),
    ]
    for i, (title, body) in enumerate(outcomes):
        col, row = i % 2, i // 2
        card(
            s,
            Inches(0.45) + col * Inches(6.4),
            Inches(1.5) + row * Inches(2.55),
            Inches(6.15),
            Inches(2.35),
            title,
            body,
        )
    footer(s, 6)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, Inches(0), Inches(0), W, H, NAVY)
    box(s, Inches(0), Inches(0), Inches(0.2), H, GOLD)
    tb(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5), "Thank you", 36, True, WHITE)
    tb(
        s,
        Inches(0.7),
        Inches(2.6),
        Inches(12),
        Inches(1.2),
        "Navy Federal Credit Union  ·  Vienna, VA\nCloud Financial Data Lakehouse",
        20,
        False,
        GOLD,
    )
    tb(
        s,
        Inches(0.7),
        Inches(4.3),
        Inches(12),
        Inches(2.0),
        "Create:   terraform apply\nLoad:     python run_local.py   or   docker compose up --build\nQuery:    Athena workgroup nfcu-cloud-analytics\nStop:     terraform destroy",
        16,
        False,
        WHITE,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
